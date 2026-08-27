"""
CSA teacher-kit deck builder.

Design system extracted directly from the Unit 1 decks that already ship
(Day1_Deck_TEACHER.pptx), so generated units are visually identical to the
pilot rather than a lookalike. Every colour, font, size and position below was
read out of that file; none of it is invented.

Slide grammar, in the order a day uses it:
    title -> warmup -> notes_preview -> objectives
    -> [section_divider -> section_content ...]
    -> worked_example -> now_break_it -> misconception
    -> vocabulary -> discussion -> end_of_day

Teacher and student editions are the same deck. The teacher edition adds
speaker notes and the answer text inside the teaching callouts; the student
edition omits both. Nothing else differs, so a teacher projecting either one
sees the same slide numbers.

No em-dashes anywhere in generated text.
"""

import math

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── palette ──────────────────────────────────────────────────────────────────
NAVY_BG     = RGBColor(0x10, 0x24, 0x3A)   # title and section-divider ground
NAVY_CHIP   = RGBColor(0x1B, 0x3A, 0x57)   # chip behind TEACHER EDITION
ACCENT      = RGBColor(0x2C, 0x6B, 0xAF)   # rules, eyebrows, key-idea labels
TINT        = RGBColor(0xF2, 0xF7, 0xFD)   # card fill
GREEN       = RGBColor(0x1E, 0x7A, 0x45)   # output / "actually true"
GREEN_TINT  = RGBColor(0xF0, 0xF7, 0xF2)
ORANGE      = RGBColor(0xC1, 0x74, 0x1C)   # misconception
ORANGE_TINT = RGBColor(0xFD, 0xF6, 0xEC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BODY        = RGBColor(0x1F, 0x29, 0x37)
MUTED       = RGBColor(0x5A, 0x6B, 0x7B)
FINEPRINT   = RGBColor(0x9A, 0xA7, 0xB4)
ON_NAVY     = RGBColor(0xCA, 0xDC, 0xFC)
ON_NAVY_DIM = RGBColor(0x9C, 0xC3, 0xE6)
TRADE_NAVY  = RGBColor(0x7F, 0x98, 0xB8)

CODE_BG     = RGBColor(0x0F, 0x22, 0x33)
CODE_TEXT   = RGBColor(0xE5, 0xED, 0xFF)
CODE_KW     = RGBColor(0x7F, 0xB4, 0xFF)
CODE_NUM    = RGBColor(0xF0, 0xC6, 0x74)
CODE_STR    = RGBColor(0x9B, 0xD1, 0x9B)
CODE_COMMENT= RGBColor(0x7F, 0x98, 0xB8)

DISPLAY = "Cambria"
UI      = "Calibri"
MONO    = "Courier New"

TRADEMARK = ("AP is a trademark of the College Board, which was not involved in "
             "the production of, and does not endorse, this resource.")

# ── geometry, in inches ──────────────────────────────────────────────────────
W, H       = 13.333, 7.5
MARGIN     = 0.50
CONTENT_W  = 12.33
FOOT_Y     = 6.94
TRADE_Y    = 7.18

# The lowest a content panel may reach before it runs into the footer.
CONTENT_BOTTOM = 6.86

# Where a card's body copy starts, measured down from the top of the card. The
# card label sits at +0.38 and the larger label needs the extra clearance.
CARD_TEXT_DY = 0.86

# Code-block sizing. Courier New is ~0.6 em wide per character, and a line of
# text occupies ~1.18 times its point size once leading is counted.
CODE_MAX_PT = 14.0
# A worked example gets the single-slide layout only when BOTH its program and
# its output fit alongside the annotations. Gating on the program alone let a
# short program with six lines of output push the OUTPUT panel through the
# footer and off the bottom of the slide, which is the same class of bug the
# wide layout was introduced to fix, one panel over.
SINGLE_SLIDE_LINES = 20
SINGLE_SLIDE_OUTPUT = 3
CODE_MIN_PT = 8.5
LINE_FACTOR = 1.18
CHAR_FACTOR = 0.60

# ── type scale, in points ────────────────────────────────────────────────────
#  Every size on every slide is named here, so "make the decks bigger" is an
#  edit to this block rather than a hunt through twelve slide builders, and the
#  hierarchy survives the edit because the relationships are visible at once.
#
#  Raised in 2026-08. The pilot was set for reading at a desk and teachers
#  projecting it could not read the smaller panels from the back of the room.
#  Nothing here went down. The panels around the text grew to match, and
#  _must_fit is what holds the two in step: a size raised past what its panel
#  can hold fails the build instead of overprinting the slide.
T_DECK_TITLE   = 42     # was 40
T_TITLE_SUB    = 18     # was 16    the deck subtitle
T_TITLE_META   = 14     # was 12    the "prepared for" line
T_TITLE_NOTE   = 14.5   # was 12.5  the which-edition-is-this line
T_TITLE_SITE   = 13     # was 11    the APCSExamPrep.com line
T_DIVIDER_NUM  = 80     # was 76
T_DIVIDER_NAME = 36     # was 34
T_DIVIDER_META = 13.5   # was 12    the divider's SECTION and topic lines
T_HEADING      = 32     # was 30

T_EYEBROW      = 12.5   # was 11
T_SUBHEAD      = 15     # was 13
T_CARD_LABEL   = 12     # was 10.5
T_CHIP         = 11.5   # was 10
T_CAPTION      = 11     # was 9
T_FOOTER       = 10.5   # was 9
T_TRADEMARK    = 9      # was 7.5

T_LEAD         = 18     # was 15    section ideas, warm-up prompt, discussion
T_BODY         = 17     # was 14.5  objectives, misconceptions, annotations
T_BODY_SM      = 15.5   # was 13    end-of-day items, vocabulary definitions
T_BODY_MIN     = 14     # was 11.5  the compact worked example's annotations
T_SECTION_NAME = 20     # was 17    notes-preview section names
T_NUMERAL      = 26     # was 24    discussion numerals
T_OUTPUT       = 13.5   # was 11.5  program output, compact layout
T_OUTPUT_WIDE  = 15.5   # was 13    program output, wide layout

# Proportional text metrics, the counterpart to CHAR_FACTOR for the code panels.
# Calibri averages nearer 0.5 em per character, but these decks are verified by
# rendering them through LibreOffice, which substitutes a wider face, so the
# estimate uses the wider number that was measured against that renderer and
# the shipped slide keeps the difference in hand as margin.
PROP_CHAR = 0.62
PROP_LINE = 1.20

JAVA_KEYWORDS = {
    'abstract','boolean','break','byte','case','catch','char','class','const',
    'continue','default','do','double','else','enum','extends','final','finally',
    'float','for','if','implements','import','instanceof','int','interface','long',
    'new','package','private','protected','public','return','short','static',
    'super','switch','this','throw','throws','try','void','while','null','true','false',
}


def _code_size(code, w, h):
    """The point size a code block renders at inside a w by h inch panel.

    A PowerPoint text frame does not clip, so code that is too long is not cut
    off: it bleeds over the slide title and past the footer. The size is
    therefore computed from the block rather than assumed, against BOTH
    constraints:

      height  lines * size * LINE_FACTOR must fit h
      width   longest line * size * CHAR_FACTOR must fit w

    Courier New is monospaced at roughly 0.6 em per character, and a line
    occupies about 1.18 times its point size once leading is included. Below
    CODE_MIN_PT the code stops being readable from the back of a room, so the
    builder raises rather than shipping a slide nobody can read: that is a
    signal to shorten the example or widen the layout.
    """
    lines = code.split('\n')
    longest = max((len(l) for l in lines), default=1)
    by_height = (h * 72.0) / (max(len(lines), 1) * LINE_FACTOR)
    by_width = (w * 72.0) / (max(longest, 1) * CHAR_FACTOR)
    size = min(CODE_MAX_PT, by_height, by_width)
    if size < CODE_MIN_PT:
        raise ValueError(
            f'code block needs {size:.1f}pt to fit ({len(lines)} lines, '
            f'longest {longest} chars) but the floor is {CODE_MIN_PT}pt. '
            'Shorten the worked example rather than shrinking the type.')
    return size


def _code_h(code, size):
    """The height in inches a code block occupies at `size` points."""
    return len(code.split('\n')) * size * LINE_FACTOR / 72.0


def _wrapped_lines(text, w, size):
    """How many lines `text` wraps to in a `w` inch column at `size` points."""
    per_line = max(1.0, (w * 72.0) / (PROP_CHAR * size))
    return max(1, math.ceil(len(text) / per_line))


def _text_h(text, w, size, line=PROP_LINE):
    """The height in inches `text` occupies in a `w` inch column once wrapped."""
    return _wrapped_lines(text, w, size) * size * line / 72.0


def _must_fit(what, text, w, size, room, line=PROP_LINE):
    """Refuse to build a slide whose text has outgrown its panel.

    A PowerPoint text frame does not clip, so text that does not fit is not cut
    off: it is drawn straight over whatever sits below it and off the bottom of
    the slide. _code has raised on exactly this since the worked-example panels
    overflowed. This is the same guard for the proportional panels, and it is
    what makes the type scale safe to raise: the failure mode is a build error
    naming the panel, not a deck that looks fine in python-pptx and is unusable
    on a projector. The fix when it fires is a bigger panel or shorter copy,
    never a smaller point size.
    """
    need = _text_h(text, w, size, line)
    if need > room:
        raise ValueError(
            f'{what}: {len(text)} characters need {need:.2f}in at {size}pt in a '
            f'{w:.2f}in column but the panel leaves {room:.2f}in. Grow the panel '
            f'or shorten the copy rather than shrinking the type.')
    return need


def _shape(slide, x, y, w, h, fill=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.shadow.inherit = False
    # python-pptx attaches a theme style to every autoshape, which is where the
    # drop shadow on text comes from. The template has none, so drop it.
    style = sh._element.find(qn('p:style'))
    if style is not None:
        sh._element.remove(style)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return sh


def _text(slide, x, y, w, h, runs, size=12, font=UI, color=BODY, bold=False,
          align=PP_ALIGN.LEFT, space_after=0, line=None, anchor=MSO_ANCHOR.TOP):
    """runs: a string, or a list of paragraphs where each paragraph is a string
    or a list of (text, color, bold) tuples."""
    sh = _shape(slide, x, y, w, h)
    tf = sh.text_frame
    tf.vertical_anchor = anchor
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after:
            p.space_after = Pt(space_after)
        if line:
            p.line_spacing = line
        pieces = para if isinstance(para, list) else [(para, color, bold)]
        for piece in pieces:
            txt, c, b = (piece if isinstance(piece, tuple) else (piece, color, bold))
            r = p.add_run()
            r.text = txt
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = b
            r.font.color.rgb = c
    return sh


def _bg(slide, color):
    """Set a solid slide background. python-pptx has no API for this, so the
    <p:bg> element is built and inserted as the first child of <p:cSld>."""
    from lxml import etree
    from pptx.oxml.ns import qn
    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    hexval = '%02X%02X%02X' % (color[0], color[1], color[2])
    xml = (
        '<p:bg xmlns:p="%s" xmlns:a="%s">'
        '<p:bgPr><a:solidFill><a:srgbClr val="%s"/></a:solidFill>'
        '<a:effectLst/></p:bgPr></p:bg>'
    ) % (P_NS, A_NS, hexval)
    bg = etree.fromstring(xml)
    cSld = slide._element.find(qn('p:cSld'))
    cSld.insert(0, bg)


class Deck:
    def __init__(self, topic, title, day, days, edition, subtitle):
        self.p = Presentation()
        self.p.slide_width = Emu(int(W * 914400))
        self.p.slide_height = Emu(int(H * 914400))
        self.topic = topic
        self.title = title
        self.day = day
        self.days = days
        self.edition = edition          # 'TEACHER' or 'STUDENT'
        self.subtitle = subtitle
        self.n = 0

    @property
    def teacher(self):
        return self.edition == 'TEACHER'

    def _new(self, navy=False):
        s = self.p.slides.add_slide(self.p.slide_layouts[6])
        self.n += 1
        if navy:
            _bg(s, NAVY_BG)
        _shape(s, 0, 0, W, 0.08, ACCENT)
        return s

    def _foot(self, s, navy=False):
        label = 'Teacher Edition' if self.teacher else 'Student Edition'
        _text(s, MARGIN, FOOT_Y, CONTENT_W, 0.24,
              f"APCSExamPrep.com  |  Topic {self.topic}  |  Day {self.day}  |  Slide {self.n}  |  {label}",
              size=T_FOOTER, color=(ON_NAVY_DIM if navy else MUTED))
        _text(s, MARGIN, TRADE_Y, CONTENT_W, 0.22, TRADEMARK,
              size=T_TRADEMARK, color=(TRADE_NAVY if navy else FINEPRINT))

    def _note(self, s, text):
        if self.teacher and text:
            s.notes_slide.notes_text_frame.text = text

    def _head(self, s, eyebrow, heading, sub=None, accent=ACCENT):
        _text(s, MARGIN, 0.30, CONTENT_W, 0.28, eyebrow.upper(), size=T_EYEBROW, bold=True, color=accent)
        # The heading measures itself and the subhead sits under whatever that
        # comes to. At the old 30pt every heading in the course was one line, so
        # a fixed subhead position was safe; at 32pt the longest section names
        # take two, and a fixed position would be written through them.
        hh = _text_h(heading, CONTENT_W, T_HEADING)
        _text(s, MARGIN, 0.62, CONTENT_W, max(0.90, hh), heading, size=T_HEADING,
              font=DISPLAY, bold=True, color=RGBColor(0x10, 0x24, 0x3A))
        if sub:
            _text(s, MARGIN, max(1.42, 0.62 + hh + 0.06), CONTENT_W, 0.46, sub,
                  size=T_SUBHEAD, color=MUTED)

    def _card(self, s, x, y, w, h, label, fill=TINT, rule=ACCENT):
        _shape(s, x, y, w, h, fill)
        _shape(s, x, y, w, 0.18, rule)
        if label:
            _text(s, x + 0.30, y + 0.38, w - 0.60, 0.30, label.upper(),
                  size=T_CARD_LABEL, bold=True, color=rule)

    # ── slide 1 ──────────────────────────────────────────────────────────────
    def title_slide(self, prepared, note=None):
        s = self._new(navy=True)
        _shape(s, 0, 0, W, 0.10, ACCENT)
        chips = [('TEACHER EDITION' if self.teacher else 'STUDENT EDITION'), f'DAY {self.day} OF {self.days}']
        for i, c in enumerate(chips):
            x = 0.60 + i * 2.50
            _shape(s, x, 0.58, 2.30, 0.40, NAVY_CHIP)
            _text(s, x, 0.58, 2.30, 0.40, c, size=T_CHIP, bold=True, color=ON_NAVY_DIM,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 0.60, 1.34, 12.10, 0.34,
              f"Unit {self.topic.split('.')[0]}  |  Topic {self.topic}  |  Day {self.day} of {self.days}",
              size=T_SUBHEAD, color=ON_NAVY_DIM)
        _must_fit('deck title', self.title, 12.10, T_DECK_TITLE, 1.74)
        _text(s, 0.60, 1.80, 12.10, 1.74, self.title, size=T_DECK_TITLE, font=DISPLAY, bold=True, color=WHITE)
        _must_fit('deck subtitle', self.subtitle, 11.40, T_TITLE_SUB, 1.00)
        _text(s, 0.60, 3.58, 11.40, 1.00, self.subtitle, size=T_TITLE_SUB, color=ON_NAVY)
        _shape(s, 0.60, 4.70, 2.40, 0.05, ACCENT)
        _text(s, 0.60, 4.96, 11.40, 0.62, prepared, size=T_TITLE_META, color=ON_NAVY_DIM)
        _text(s, 0.60, 5.68, 11.40, 0.62,
              ("Teacher edition. Speaker notes carry the answers and the misconceptions for each slide."
               if self.teacher else
               "Student edition. Fill in your guided notes as we go."),
              size=T_TITLE_NOTE, color=ON_NAVY)
        _text(s, 0.60, 6.50, 5.00, 0.32, "APCSExamPrep.com", size=T_TITLE_SITE, bold=True, color=WHITE)
        _text(s, 0.60, 6.94, 12.10, 0.40, TRADEMARK, size=T_TRADEMARK, color=TRADE_NAVY)
        self._note(s, note)

    # ── warm-up ──────────────────────────────────────────────────────────────
    def warmup(self, heading, prompt, draw_out=None):
        s = self._new()
        self._head(s, 'WARM-UP', heading)
        w = CONTENT_W - 0.60
        self._card(s, MARGIN, 1.90, CONTENT_W, 2.32, 'ON THE BOARD')
        h = _must_fit('warm-up prompt', prompt, w, T_LEAD, 2.32 - CARD_TEXT_DY - 0.14, 1.25)
        _text(s, MARGIN + 0.30, 1.90 + CARD_TEXT_DY, w, h, prompt, size=T_LEAD, color=BODY,
              space_after=6, line=1.25)
        if self.teacher and draw_out:
            self._card(s, MARGIN, 4.38, CONTENT_W, 2.34, 'WORTH DRAWING OUT', GREEN_TINT, GREEN)
            h = _must_fit('warm-up draw-out', draw_out, w, T_BODY, 2.34 - CARD_TEXT_DY - 0.14, 1.2)
            _text(s, MARGIN + 0.30, 4.38 + CARD_TEXT_DY, w, h, draw_out, size=T_BODY,
                  color=BODY, space_after=5, line=1.2)
        self._foot(s)
        self._note(s, draw_out)

    # ── guided notes preview ─────────────────────────────────────────────────
    def notes_preview(self, sections):
        s = self._new()
        self._head(s, 'GUIDED NOTES', 'What you will be filling in',
                   "The sections of today's guided notes packet, in the order they appear.")
        n = max(1, len(sections))
        w = (CONTENT_W - 0.30 * (n - 1)) / n
        cw = w - 0.56
        for i, (name, blurb) in enumerate(sections):
            x = MARGIN + i * (w + 0.30)
            self._card(s, x, 2.10, w, 4.20, f'SECTION {i + 1}')
            nh = max(0.62, _must_fit('notes-preview name', name, cw, T_SECTION_NAME, 1.20))
            _text(s, x + 0.28, 2.10 + CARD_TEXT_DY, cw, nh, name, size=T_SECTION_NAME,
                  font=DISPLAY, bold=True, color=RGBColor(0x10, 0x24, 0x3A))
            by = 2.10 + CARD_TEXT_DY + nh + 0.18
            bh = _must_fit('notes-preview blurb', blurb, cw, T_BODY, 6.30 - 0.16 - by)
            _text(s, x + 0.28, by, cw, bh, blurb, size=T_BODY, color=BODY, line=1.2)
        self._foot(s)
        self._note(s, "The sections of today's guided notes packet, in the order they appear.")

    # ── objectives ───────────────────────────────────────────────────────────
    def objectives(self, items):
        s = self._new()
        self._head(s, 'LESSON OBJECTIVES', 'By the end of today you will be able to',
                   'The CED learning objectives for this topic, with their codes.')
        # This heading is fixed text and always sets one line, so the cards can
        # start directly under the subhead rather than at a worst-case offset.
        y = 1.80
        for text, code in items:
            self._card(s, MARGIN, y, CONTENT_W, 1.58, 'CB REQUIRED')
            h = _must_fit('objective', text, CONTENT_W - 2.20, T_BODY, 1.58 - CARD_TEXT_DY - 0.12)
            _text(s, MARGIN + 0.30, y + CARD_TEXT_DY, CONTENT_W - 2.20, h, text, size=T_BODY, color=BODY)
            _text(s, W - MARGIN - 1.80, y + 0.38, 1.50, 0.32, code, size=T_CARD_LABEL, bold=True,
                  color=ACCENT, align=PP_ALIGN.RIGHT)
            y += 1.72
        self._foot(s)
        self._note(s, 'The CED learning objectives for this topic, with their codes.')

    # ── section divider ──────────────────────────────────────────────────────
    def section_divider(self, index, name):
        s = self._new(navy=True)
        _text(s, 0.60, 2.14, 3.00, 1.50, f"{index:02d}", size=T_DIVIDER_NUM, font=DISPLAY,
              bold=True, color=ACCENT)
        _text(s, 0.60, 3.64, 11.40, 0.36, 'SECTION', size=T_DIVIDER_META, bold=True, color=ON_NAVY_DIM)
        _must_fit('divider name', name, 11.40, T_DIVIDER_NAME, 1.44)
        _text(s, 0.60, 4.02, 11.40, 1.44, name, size=T_DIVIDER_NAME, font=DISPLAY, bold=True, color=WHITE)
        _text(s, 0.60, 5.58, 11.40, 0.36, f"Topic {self.topic}  |  Day {self.day}",
              size=T_DIVIDER_META, color=ON_NAVY_DIM)
        self._foot(s, navy=True)

    # ── section content ──────────────────────────────────────────────────────
    def section_content(self, index, name, ideas, part=None, note=None):
        s = self._new()
        heading = name if not part else f"{name} ({part})"
        self._head(s, f'SECTION {index:02d}', heading, f'Guided notes, section {index}.')
        card_y, card_h = 2.10, 4.36
        self._card(s, MARGIN, card_y, CONTENT_W, card_h, 'KEY IDEA')
        note_y = card_y + card_h - 0.44
        w = CONTENT_W - 0.60
        y = card_y + CARD_TEXT_DY
        for idea in ideas:
            # Each idea measures itself rather than guessing a height from a
            # character count, so the stack below it cannot be pushed off the card.
            h = _must_fit('section idea', idea, w, T_LEAD, note_y - 0.12 - y, 1.22)
            _text(s, MARGIN + 0.30, y, w, h, idea, size=T_LEAD, color=BODY, line=1.22)
            y += h + 0.24
        _text(s, MARGIN + 0.30, note_y, w, 0.30,
              f"Topic {self.topic}  |  Section {index}", size=T_CAPTION, color=MUTED)
        self._foot(s)
        self._note(s, note or 'Where students go quiet here it is usually the vocabulary rather than the concept.')

    # ── worked example ───────────────────────────────────────────────────────
    def worked_example(self, heading, code, notice, output, caption='Complete and runnable as shown.', note=None):
        """One slide for a short program, two for a long one.

        The pilot deck's example was a twelve line Greeting class, and the
        panel was sized for it. Real worked examples in this course are whole
        classes: constructors, getters and a main, routinely thirty lines and
        occasionally fifty. Rendering those into the original panel pushed the
        code over the slide title and past the footer, because a PowerPoint
        text frame does not clip.

        Shrinking the type was not the answer either. Nobody reads 6pt Java
        from the back of a room. So a long program gets the full slide width in
        two columns, and its annotations move to a second slide the teacher
        flips to. Short programs keep the original single-slide layout, which
        is better when it fits.

        "Fits" now means the output as well as the program. The single-slide
        OUTPUT panel is the shortest thing on the slide and sits directly above
        the footer, so a twelve line program printing six lines of output was
        the same overflow one panel over. Either half being long sends the
        example wide, where each half has a slide to itself.
        """
        lines = code.split('\n')
        if len(lines) <= SINGLE_SLIDE_LINES and len(output) <= SINGLE_SLIDE_OUTPUT:
            self._worked_compact(heading, code, notice, output, caption, note)
        else:
            self._worked_wide(heading, code, notice, output, caption, note)

    def _worked_compact(self, heading, code, notice, output, caption, note):
        s = self._new()
        self._head(s, 'WORKED EXAMPLE', heading)
        # The panels start straight under the heading and run to the footer, so
        # the code gets every inch of height available: the compact layout is
        # height-bound, never width-bound, at every program length it accepts.
        # The panel is drawn to the height the code turns out to need, so a
        # twelve line program does not sit in the middle of a dark box sized
        # for a twenty line one.
        size = _code_size(code, 6.67, 3.40)
        ch = _code_h(code, size)
        panel_h = ch + 0.26
        card_h = min(5.10, 0.92 + panel_h + 0.44)
        self._card(s, MARGIN, 1.74, 7.55, card_h, 'THE COMPLETE PROGRAM')
        _shape(s, 0.80, 2.66, 6.95, panel_h, CODE_BG)
        self._code(s, 0.94, 2.79, 6.67, ch, code, size)
        _text(s, 0.80, 1.74 + card_h - 0.34, 6.95, 0.28, caption, size=T_CAPTION, color=MUTED)

        self._card(s, 8.38, 1.74, 4.40, 3.26, 'WHAT TO NOTICE')
        y = 1.74 + CARD_TEXT_DY
        for item in notice:
            h = _must_fit('worked-example note', item, 3.62, T_BODY_MIN, 5.00 - 0.12 - y, 1.12)
            _text(s, 8.62, y, 0.20, h, '\u2022', size=T_BODY_MIN, color=ACCENT)
            _text(s, 8.88, y, 3.62, h, item, size=T_BODY_MIN, color=BODY, line=1.12)
            y += h + 0.13

        self._card(s, 8.38, 5.18, 4.40, 1.66, 'OUTPUT', GREEN_TINT, GREEN)
        y = 5.18 + CARD_TEXT_DY
        for line in output:
            _text(s, 8.62, y, 3.88, 0.28, line, size=T_OUTPUT, font=MONO, color=BODY)
            y += T_OUTPUT * 1.15 / 72.0
        self._foot(s)
        self._note(s, note or 'A complete, runnable program. The annotations are on the next slide.')

    def _worked_wide(self, heading, code, notice, output, caption, note):
        # Slide one: the program, full width, two columns.
        s = self._new()
        self._head(s, 'WORKED EXAMPLE', heading)
        lines = code.split('\n')
        half = (len(lines) + 1) // 2
        left, right = '\n'.join(lines[:half]), '\n'.join(lines[half:])
        colw = 5.72
        # Both columns render at one size, and that size is what BOTH of them
        # can hold. Sizing the left column and pinning the right one to it
        # assumed the left column had the longest line, and where it did not
        # the right column was drawn off the right-hand edge of the slide.
        size = min(_code_size(left, colw, 3.68), _code_size(right, colw, 3.68))
        ch = max(_code_h(left, size), _code_h(right, size))
        panel_h = ch + 0.26
        card_h = min(5.10, 0.84 + panel_h + 0.44)
        self._card(s, MARGIN, 1.74, CONTENT_W, card_h, 'THE COMPLETE PROGRAM')
        # The card label occupies 2.12 to 2.32, so the panel starts below it.
        _shape(s, 0.72, 2.58, 11.89, panel_h, CODE_BG)
        self._code(s, 0.90, 2.71, colw, ch, left, size)
        self._code(s, 6.86, 2.71, colw, ch, right, size)
        _text(s, 0.72, 1.74 + card_h - 0.34, CONTENT_W, 0.28,
              caption + '  Continues in the right column.', size=T_CAPTION, color=MUTED)
        self._foot(s)
        self._note(s, note or 'A complete, runnable program. The annotations are on the next slide.')

        # Slide two: what to notice, and the output. Both panels have a whole
        # slide here rather than a column, which is what lets this layout carry
        # the longest annotations and the longest output in the course.
        s2 = self._new()
        self._head(s2, 'WORKED EXAMPLE', heading, 'What to notice in the program on the previous slide.')
        self._card(s2, MARGIN, 2.10, 7.55, 4.20, 'WHAT TO NOTICE')
        y = 2.10 + CARD_TEXT_DY
        for item in notice:
            h = _must_fit('worked-example note', item, 6.45, T_BODY, 6.30 - 0.12 - y, 1.15)
            _text(s2, 0.80, y, 0.22, h, '\u2022', size=T_BODY, color=ACCENT)
            _text(s2, 1.10, y, 6.45, h, item, size=T_BODY, color=BODY, line=1.15)
            y += h + 0.16
        self._card(s2, 8.38, 2.10, 4.40, 4.20, 'OUTPUT', GREEN_TINT, GREEN)
        y = 2.10 + CARD_TEXT_DY
        for line in output:
            _text(s2, 8.62, y, 3.88, 0.30, line, size=T_OUTPUT_WIDE, font=MONO, color=BODY)
            y += T_OUTPUT_WIDE * 1.15 / 72.0
        self._foot(s2)
        self._note(s2, 'The output the program on the previous slide produces.')

    def _code(self, s, x, y, w, h, code, size):
        """Render a code block at an already-decided point size.

        The size comes from _code_size rather than from here, because both
        columns of a wide worked example have to agree on one size and the
        panel behind them has to be drawn to the height the code will actually
        occupy. All three need the number before any of them can be drawn.
        """
        lines = code.split('\n')
        sh = _shape(s, x, y, w, h)
        tf = sh.text_frame
        tf.word_wrap = False
        for i, raw in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.0
            p.space_after = Pt(0)
            for txt, col in _tokenize(raw):
                r = p.add_run()
                r.text = txt
                r.font.name = MONO
                r.font.size = Pt(size)
                r.font.color.rgb = col
        return size

    # ── now break it ─────────────────────────────────────────────────────────
    def now_break_it(self, change, happens, why, note=None):
        s = self._new()
        self._head(s, 'NOW BREAK IT  |  ONE CHANGE', 'Change one line and run it again')
        room = 2.80 - CARD_TEXT_DY - 0.12
        self._card(s, MARGIN, 1.40, 6.00, 2.80, 'THE CHANGE')
        h = _must_fit('break-it change', change, 5.40, T_BODY, room, 1.2)
        _text(s, MARGIN + 0.30, 1.40 + CARD_TEXT_DY, 5.40, h, change, size=T_BODY, color=BODY, line=1.2)
        self._card(s, 6.83, 1.40, 6.00, 2.80, 'WHAT HAPPENS', ORANGE_TINT, ORANGE)
        h = _must_fit('break-it outcome', happens, 5.40, T_BODY, room, 1.2)
        _text(s, 7.13, 1.40 + CARD_TEXT_DY, 5.40, h, happens, size=T_BODY, color=BODY, line=1.2)
        self._card(s, MARGIN, 4.38, CONTENT_W, 2.44, 'WHY THIS MATTERS', GREEN_TINT, GREEN)
        h = _must_fit('break-it rationale', why, CONTENT_W - 0.60, T_BODY, 2.44 - CARD_TEXT_DY - 0.12, 1.2)
        _text(s, MARGIN + 0.30, 4.38 + CARD_TEXT_DY, CONTENT_W - 0.60, h, why, size=T_BODY,
              color=BODY, line=1.2)
        self._foot(s)
        self._note(s, note or 'One change to the program on the previous slide. This is the error pattern the exam tests on this topic.')

    # ── misconception ────────────────────────────────────────────────────────
    def misconception(self, heading, think, truth, note=None):
        s = self._new()
        self._head(s, 'COMMON MISCONCEPTION', heading, accent=ORANGE)
        # This heading runs to two lines at the longest, so the cards clear it.
        card_y, card_h = 1.80, 4.60
        room = card_h - CARD_TEXT_DY - 0.12
        self._card(s, MARGIN, card_y, 6.00, card_h, 'WHAT STUDENTS THINK', ORANGE_TINT, ORANGE)
        h = _must_fit('misconception belief', think, 5.40, T_BODY, room, 1.22)
        _text(s, MARGIN + 0.30, card_y + CARD_TEXT_DY, 5.40, h, think, size=T_BODY, color=BODY, line=1.22)
        self._card(s, 6.83, card_y, 6.00, card_h, 'WHAT IS ACTUALLY TRUE', GREEN_TINT, GREEN)
        h = _must_fit('misconception correction', truth, 5.40, T_BODY, room, 1.22)
        _text(s, 7.13, card_y + CARD_TEXT_DY, 5.40, h, truth, size=T_BODY, color=BODY, line=1.22)
        _text(s, MARGIN, 6.50, CONTENT_W, 0.30, f"Topic {self.topic}", size=T_CAPTION, color=MUTED)
        self._foot(s)
        self._note(s, note or 'This is the misconception that costs the most marks on this topic.')

    # ── vocabulary ───────────────────────────────────────────────────────────
    def vocabulary(self, terms):
        s = self._new()
        self._head(s, 'KEY VOCABULARY', 'Words you need to use precisely',
                   'The terms the exam uses in its question stems.')
        cols, colw = 3, (CONTENT_W - 0.60) / 3
        card_h = 2.30
        for i, (term, definition) in enumerate(terms[:6]):
            cx = MARGIN + (i % cols) * (colw + 0.30)
            cy = 1.78 + (i // cols) * (card_h + 0.16)
            self._card(s, cx, cy, colw, card_h, term)
            h = _must_fit('vocabulary definition', definition, colw - 0.56, T_BODY_SM,
                          card_h - CARD_TEXT_DY - 0.12, 1.15)
            _text(s, cx + 0.28, cy + CARD_TEXT_DY, colw - 0.56, h, definition,
                  size=T_BODY_SM, color=BODY, line=1.15)
        self._foot(s)
        self._note(s, 'The vocabulary for this topic as the CED uses it.')

    # ── discussion ───────────────────────────────────────────────────────────
    def discussion(self, questions, note=None):
        s = self._new()
        self._head(s, 'TALK IT THROUGH', 'Discussion')
        y, card_h = 1.90, 2.06
        for i, q in enumerate(questions, 1):
            self._card(s, MARGIN, y, CONTENT_W, card_h, None)
            _text(s, MARGIN + 0.30, y + 0.42, 0.60, 0.60, str(i), size=T_NUMERAL, font=DISPLAY,
                  bold=True, color=ACCENT)
            h = _must_fit('discussion question', q, CONTENT_W - 1.40, T_LEAD, card_h - 0.52 - 0.14, 1.2)
            _text(s, MARGIN + 1.00, y + 0.52, CONTENT_W - 1.40, h, q, size=T_LEAD, color=BODY, line=1.2)
            y += card_h + 0.18
        self._foot(s)
        self._note(s, note)

    # ── end of day ───────────────────────────────────────────────────────────
    def end_of_day(self, learned, up_next, extra, note=None):
        s = self._new()
        self._head(s, f'END OF DAY {self.day}  |  DAY {self.day} OF {self.days}', f'End of day {self.day}')
        self._card(s, MARGIN, 1.80, 7.30, 4.38, 'TODAY YOU LEARNED')
        y = 1.80 + CARD_TEXT_DY
        for item in learned:
            h = _must_fit('end-of-day outcome', item, 6.70, T_BODY, 6.18 - 0.12 - y, 1.18)
            _text(s, MARGIN + 0.30, y, 6.70, h, item, size=T_BODY, color=BODY, line=1.18)
            y += h + 0.22
        self._card(s, 8.10, 1.80, 4.73, 2.10, 'UP NEXT')
        h = _must_fit('end-of-day up next', up_next, 4.13, T_BODY_SM, 2.10 - CARD_TEXT_DY - 0.12, 1.2)
        _text(s, 8.40, 1.80 + CARD_TEXT_DY, 4.13, h, up_next, size=T_BODY_SM, color=BODY, line=1.2)
        self._card(s, 8.10, 4.08, 4.73, 2.10, 'EXTRA PRACTICE', GREEN_TINT, GREEN)
        h = _must_fit('end-of-day extra practice', extra, 4.13, T_BODY_SM, 2.10 - CARD_TEXT_DY - 0.12, 1.2)
        _text(s, 8.40, 4.08 + CARD_TEXT_DY, 4.13, h, extra, size=T_BODY_SM, color=BODY, line=1.2)
        self._foot(s)
        self._note(s, note or 'The objectives restated as outcomes, plus what the next day covers.')

    def save(self, path):
        self.p.save(path)


def _tokenize(line):
    """Very small Java tokenizer, matching the highlight colours the Unit 1
    decks use: keywords blue, numbers amber, strings green, comments grey."""
    out, i, n = [], 0, len(line)
    stripped = line.lstrip()
    if stripped.startswith('//'):
        return [(line, CODE_COMMENT)]
    buf = ''
    while i < n:
        ch = line[i]
        if ch == '"':
            if buf:
                out.append((buf, CODE_TEXT)); buf = ''
            j = i + 1
            while j < n and line[j] != '"':
                j += 1
            out.append((line[i:min(j + 1, n)], CODE_STR))
            i = j + 1
            continue
        if ch.isalpha() or ch == '_':
            j = i
            while j < n and (line[j].isalnum() or line[j] == '_'):
                j += 1
            word = line[i:j]
            if word in JAVA_KEYWORDS:
                if buf:
                    out.append((buf, CODE_TEXT)); buf = ''
                out.append((word, CODE_KW))
            else:
                buf += word
            i = j
            continue
        if ch.isdigit():
            j = i
            while j < n and (line[j].isdigit() or line[j] == '.'):
                j += 1
            if buf:
                out.append((buf, CODE_TEXT)); buf = ''
            out.append((line[i:j], CODE_NUM))
            i = j
            continue
        buf += ch
        i += 1
    if buf:
        out.append((buf, CODE_TEXT))
    return out or [(' ', CODE_TEXT)]
