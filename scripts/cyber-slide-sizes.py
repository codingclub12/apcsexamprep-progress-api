#!/usr/bin/env python3
"""
Read the real font sizes out of exported Teacher Bundle decks.

WHY THIS EXISTS. scripts/slide-type-bump.gs has a preview() that reports the
same histogram, and it is the only one of the two that can run against Drive
directly. But it runs inside Tanner's Apps Script project, dies at six minutes,
and walks the deck table in order: the first real preview() run reached 136 of
224 AP CSP decks and never opened a single cyber one, while printing a single
histogram that looked like it described both courses.

This reads exported decks instead, which has none of those limits. It produced
the complete AP CYBER census on 2026-09-02: all 70 decks, 16,900 runs, nothing
sampled. It cannot write anything, so it does not replace the Apps Script; it
replaces the guessing that would otherwise go into a ladder.

GETTING THE DECKS. Export each Slides file as .pptx. Either File > Download in
Slides, or the Drive connector with
exportMimeType=application/vnd.openxmlformats-officedocument.presentationml.presentation,
which is how the 2026-09-02 census was taken. The file ids are in
config/csp-slide-embeds.js and config/cyber-slide-embeds.js, reachable through
their slideId() with the matching manifest.

Usage:
    python3 scripts/cyber-slide-sizes.py <dir-of-pptx>
    python3 scripts/cyber-slide-sizes.py <dir> --floor 10 --ceiling 18

Needs python-pptx.

THE INHERITANCE PROBLEM, which is the whole reason this is not ten lines.
python-pptx reports None for a run whose size is not set explicitly. Slides'
getFontSize() does not: it resolves the chain and returns the size the text is
actually drawn at. A naive read undercounts badly. On the first cyber deck, 60
of 236 runs came back None and every one resolved to 14pt, which is the largest
bucket in the corpus. So this walks the same chain PowerPoint does: run,
paragraph, the shape's placeholder in the layout, then in the master, then the
master's txStyles for that placeholder type, then the presentation default.
Across all 70 cyber decks it resolved 16,900 of 16,900.
"""

import argparse
import glob
import math
import os
import sys
from collections import Counter

from pptx import Presentation
from pptx.oxml.ns import qn

def _lvl(para):
    try:
        return int(para._pPr.get('lvl')) if para._pPr is not None and para._pPr.get('lvl') else 0
    except Exception:
        return 0

def _sz_from_lststyle(lst, lvl):
    if lst is None:
        return None
    tag = qn('a:lvl%dpPr' % (lvl + 1))
    p = lst.find(tag)
    if p is None:
        return None
    d = p.find(qn('a:defRPr'))
    if d is None or d.get('sz') is None:
        return None
    return int(d.get('sz')) / 100.0

def _ph_type(shape):
    try:
        return shape.placeholder_format.type, shape.placeholder_format.idx
    except Exception:
        return None, None

def _match_ph(container, ptype, idx):
    if container is None:
        return None
    for sh in container.placeholders:
        try:
            if sh.placeholder_format.idx == idx:
                return sh
        except Exception:
            continue
    for sh in container.placeholders:
        try:
            if sh.placeholder_format.type == ptype:
                return sh
        except Exception:
            continue
    return None

def _txstyle_for(master, ptype):
    txs = master._element.find(qn('p:txStyles'))
    if txs is None:
        return None
    name = 'p:otherStyle'
    s = str(ptype)
    if 'TITLE' in s or 'CENTER_TITLE' in s:
        name = 'p:titleStyle'
    elif 'BODY' in s or 'SUBTITLE' in s or 'OBJECT' in s or 'None' in s:
        name = 'p:bodyStyle'
    return txs.find(qn(name))

def effective_sizes(prs):
    """Yield the effective point size of every run in the presentation."""
    out = []
    for slide in prs.slides:
        layout = slide.slide_layout
        master = layout.slide_master if layout is not None else None
        for shape in slide.shapes:
            for sz in _shape_sizes(shape, layout, master, prs):
                out.append(sz)
    return out

def _shape_sizes(shape, layout, master, prs):
    res = []
    if shape.shape_type is not None and str(shape.shape_type).startswith('GROUP'):
        for sub in shape.shapes:
            res.extend(_shape_sizes(sub, layout, master, prs))
        return res
    if getattr(shape, 'has_table', False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for r in para.runs:
                        res.append(r.font.size.pt if r.font.size else None)
        return res
    if not getattr(shape, 'has_text_frame', False) or not shape.has_text_frame:
        return res

    ptype, idx = _ph_type(shape)
    lay_ph = _match_ph(layout, ptype, idx) if layout is not None else None
    mas_ph = _match_ph(master, ptype, idx) if master is not None else None

    for para in shape.text_frame.paragraphs:
        lvl = _lvl(para)
        chain = []
        if para.font.size:
            chain.append(para.font.size.pt)
        for src in (shape, lay_ph, mas_ph):
            if src is None:
                continue
            lst = src.text_frame._txBody.find(qn('a:lstStyle'))
            chain.append(_sz_from_lststyle(lst, lvl))
        if master is not None:
            chain.append(_sz_from_lststyle(_txstyle_for(master, ptype), lvl))
        dts = prs._element.find(qn('p:defaultTextStyle'))
        chain.append(_sz_from_lststyle(dts, lvl))
        inherited = next((c for c in chain if c), None)
        for r in para.runs:
            res.append(r.font.size.pt if r.font.size else inherited)
    return res


# ── the ladder, mirroring proposeLadder_ in scripts/slide-type-bump.gs ───────
# Kept in step with that function on purpose: a ladder proposed here and a
# ladder proposed there must agree, or the number pasted into LADDERS would not
# be the one the script would have chosen.

def build_ladder(sizes, lift, floor, ceiling):
    out, prev = [], None
    for s in sizes:
        raw = s + lift * (ceiling - s) / (ceiling - floor)
        t = round(raw * 2) / 2
        if t < s:
            t = s
        if prev is not None and t <= prev:
            t = prev + 0.5
        if t >= ceiling:
            return None
        out.append((s, t))
        prev = t
    return out


def propose_ladder(sizes, floor, ceiling, max_lift=2.5):
    in_range = sorted(s for s in sizes if floor <= s < ceiling)
    for k in range(int(round(max_lift * 4)), -1, -1):
        out = build_ladder(in_range, k / 4, floor, ceiling)
        if out is not None:
            return out, k / 4
    return [], 0.0


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('directory')
    ap.add_argument('--floor', type=float, default=10)
    ap.add_argument('--ceiling', type=float, default=18)
    ap.add_argument('--course', default='ap-cybersecurity',
                    help='only used to label the paste-ready block')
    args = ap.parse_args()

    files = sorted(glob.glob(os.path.join(args.directory, '*.pptx')))
    if not files:
        sys.exit(f'no .pptx files in {args.directory}')

    hist = Counter()
    decks = unresolved = total = 0
    for f in files:
        try:
            sizes = effective_sizes(Presentation(f))
        except Exception as exc:
            print(f'skip {os.path.basename(f)}: {exc}')
            continue
        decks += 1
        for s in sizes:
            total += 1
            if s is None:
                unresolved += 1
            else:
                hist[float(s)] += 1

    in_range = sum(v for k, v in hist.items() if args.floor <= k < args.ceiling)
    print(f'decks read      : {decks}')
    print(f'text runs       : {total}   unresolved: {unresolved}')
    if unresolved:
        print('  WARNING: unresolved runs are runs whose size this script could not')
        print('  work out. The Apps Script would see a number there, so the histogram')
        print('  below is incomplete and the ladder it proposes may be wrong.')
    print(f'runs in {args.floor:g} to {args.ceiling:g}pt : {in_range}'
          f'  ({in_range * 100 // max(total, 1)}%)')
    print()
    print('size histogram (pt: runs, per deck):')
    for k in sorted(hist):
        print(f'  {k:>6g}: {hist[k]:>6}   {hist[k] / max(decks, 1):>6.1f}/deck')

    ladder, lift = propose_ladder(sorted(hist), args.floor, args.ceiling)
    print()
    print(f'proposed ladder (lift {lift:g}pt at the floor, tapering to nothing '
          f'at {args.ceiling:g}):')
    band = [s for s in hist if 10 <= s <= 14]
    lut = dict(ladder)
    if band and all(s in lut for s in band):
        mean = (sum(hist[s] * (lut[s] - s) for s in band)
                / sum(hist[s] for s in band))
        print(f'  mean lift across the original 10 to 14 band: {mean:.2f}pt')
    moved = sum(hist[s] for s, _ in ladder)
    print(f'  runs moved: {moved} of {total}')

    # Order is the property that matters, so check it here too rather than
    # trusting that the generator got it right.
    after = [(s, lut.get(s, s)) for s in sorted(hist)]
    broken = sum(hist[a[0]] for i, a in enumerate(after)
                 for b in after[i + 1:] if a[1] >= b[1])
    landings = [t for _, t in after]
    print(f'  runs with order broken: {broken}')
    print(f'  every size lands somewhere distinct: '
          f'{len(set(landings)) == len(landings)}')

    print()
    print('paste into LADDERS in scripts/slide-type-bump.gs:')
    print(f"    '{args.course}': {{")
    for s, t in ladder:
        print(f"      '{s:g}': {t:g},")
    print('    },')


if __name__ == '__main__':
    main()
