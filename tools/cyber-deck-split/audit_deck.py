#!/usr/bin/env python3
"""Report the per-day edits a freshly split deck still needs.

    python3 audit_deck.py <deck.pptx> --day N --of M

`split_deck.py` does the range cut and nothing else. Everything it leaves
undone is listed in docs/cyber-unit3-tier1-split-spec.md under "Edits that
apply to every day deck", and every item there is visible to a teacher: a Day 2
deck whose title slide still says DAY 1 OF 1, whose footers count to the whole
lesson's slide total, and whose speaker notes open by announcing it covers both
learning objectives.

This exists because a split deck LOOKS finished. It opens, it paginates, every
slide renders. The spec's own warning is that a pure range cut can leave a day
with no practice and no close, and nothing about the file says so.

Exit code is 1 while anything is outstanding, so this can gate an upload.
"""
import argparse
import re
import sys
from pptx import Presentation

FOOTER = re.compile(r'Slide\s+(\d+)\s+of\s+(\d+)', re.I)
BADGE = re.compile(r'DAY\s+(\d+)\s+OF\s+(\d+)', re.I)
WHOLE_LESSON = re.compile(
    r'single comprehensive|both learning objectives|covering both|whole lesson',
    re.I)


def slide_text(slide):
    return '\n'.join(sh.text_frame.text for sh in slide.shapes
                     if sh.has_text_frame)


def audit(path, day, of):
    prs = Presentation(path)
    n = len(prs.slides)
    problems = []

    # 1. Footer denominators are per-day in this course and restart each day.
    wrong = []
    for i, s in enumerate(prs.slides, 1):
        for _, denom in FOOTER.findall(slide_text(s)):
            if int(denom) != n:
                wrong.append((i, denom))
    if wrong:
        shown = ', '.join(f'slide {i} reads "of {d}"' for i, d in wrong[:4])
        problems.append(
            f'footer denominator: {len(wrong)} slide(s) do not read "of {n}" '
            f'({shown}{", ..." if len(wrong) > 4 else ""})')

    # 2. Title badge must name this day and the real day count.
    title = slide_text(prs.slides[0])
    badges = BADGE.findall(title)
    if not badges:
        problems.append('title slide has no DAY n OF m badge')
    else:
        for a, b in badges:
            if (int(a), int(b)) != (day, of):
                problems.append(
                    f'title badge reads DAY {a} OF {b}, expected DAY {day} OF {of}')

    # 3. Title-slide speaker notes narrate the whole lesson.
    if prs.slides[0].has_notes_slide:
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        if WHOLE_LESSON.search(notes):
            m = WHOLE_LESSON.search(notes)
            problems.append(
                f'title speaker notes still narrate the whole lesson '
                f'("...{notes[max(0, m.start()-25):m.end()+25].strip()}...")')

    # 4. A teaching day that ends with no practice and no close. The spec calls
    #    this out for 3.4 Day 1 specifically; it is worth checking everywhere.
    tail = slide_text(prs.slides[-1]).upper()
    if not any(k in tail for k in
               ('STOP AND THINK', 'END OF', 'TODAY YOU LEARNED', 'IN ONE SLIDE',
                'YOUR TURN')):
        first = tail.strip().splitlines()[0][:60] if tail.strip() else '(no text)'
        problems.append(
            f'last slide is neither practice nor a close (reads "{first}")')

    # 5. Guided-notes preview must name this day.
    for i, s in enumerate(prs.slides, 1):
        t = slide_text(s)
        if 'GUIDED NOTES' in t.upper() and 'What You' in t:
            m = re.search(r'Your Day (\d+) Guided Notes', t)
            if m and int(m.group(1)) != day:
                problems.append(
                    f'slide {i} guided-notes preview says "Your Day {m.group(1)}"'
                    f', expected Day {day}')
    # The badge and footer can appear in more than one shape on a slide, so the
    # same fact can be discovered twice. Report each distinct problem once.
    seen, unique = set(), []
    for p in problems:
        if p not in seen:
            seen.add(p)
            unique.append(p)
    return n, unique


if __name__ == '__main__':
    ap = argparse.ArgumentParser()
    ap.add_argument('deck')
    ap.add_argument('--day', type=int, required=True)
    ap.add_argument('--of', type=int, required=True)
    a = ap.parse_args()
    n, problems = audit(a.deck, a.day, a.of)
    label = f'{a.deck}  ({n} slides, day {a.day} of {a.of})'
    if not problems:
        print(f'OK   {label}')
        sys.exit(0)
    print(f'TODO {label}')
    for p in problems:
        print(f'  - {p}')
    sys.exit(1)
