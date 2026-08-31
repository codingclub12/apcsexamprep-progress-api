#!/usr/bin/env python3
"""Author the per-day decks for CED 3.3 Segmentation and 3.4 Firewalls.

    python3 build_33_34.py <workdir>

Consumes <workdir>/<lesson>_src_<variant>.pptx and <workdir>/cut/*.pptx and
writes finished decks to <workdir>/final/.

Cut lists come from docs/cyber-unit3-tier1-split-spec.md and were re-verified
against the shipped decks before use: 3.3 matched 20 of 20, 3.4 matched 21 of
21. The spec's one wrong claim is that "the STUDENT deck of each pair has the
same structure and takes the same cuts". It does not. Each student deck drops
its lesson's teacher-only Enrichment slide, so it is one slide shorter and its
Day 2 ends one earlier. Both enrichment slides fall after the seam, so the seam
itself is unchanged.

Shapes are located by geometry and text, never by index, because the variants
differ: 3.3 STUDENT's guided-notes slide has two section cards where TEACHER
has three, and the TEACHER title slide carries an extra note shape.

EK CODES. CLAUDE.md forbids CED Essential Knowledge codes in front of students.
The shipped decks break that rule heavily and cleaning them is a separate,
tracked job. Slides authored HERE follow the rule: teacher decks cite EKs,
student decks do not. That leaves a new student slide sitting next to old ones
that still carry codes. The alternative was to add more of exactly what has
already been flagged for removal.
"""
import copy
import sys
import os
from pptx import Presentation
from pptx.util import Emu, Inches

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import deckedit as de

IN = lambda emu: Emu(emu).inches


# --------------------------------------------------------------------------
# Locating things without relying on shape order
# --------------------------------------------------------------------------
def find(slide, needle, start=False):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if (t.startswith(needle) if start else needle in t):
            return sh
    raise KeyError(f'no shape containing {needle!r}')


def columns(slide, ytop, ybot):
    """Group the guided-notes cards into columns by their left edge."""
    buckets = {}
    for sh in slide.shapes:
        y = IN(sh.top)
        if not (ytop <= y <= ybot):
            continue
        buckets.setdefault(round(IN(sh.left), 1), []).append(sh)
    return [buckets[k] for k in sorted(buckets)]


def rows(slide, ytop, ybot):
    """Group the objective blocks into rows by their top edge."""
    buckets = {}
    for sh in slide.shapes:
        y = IN(sh.top)
        if not (ytop <= y <= ybot):
            continue
        placed = False
        for k in list(buckets):
            if abs(k - y) < 0.35:
                buckets[k].append(sh)
                placed = True
                break
        if not placed:
            buckets[y] = [sh]
    return [buckets[k] for k in sorted(buckets)]


def drop_group(group):
    for sh in group:
        sh._element.getparent().remove(sh._element)


def shift(group, dx=0.0, dy=0.0):
    for sh in group:
        sh.left = Emu(sh.left + Inches(dx))
        sh.top = Emu(sh.top + Inches(dy))


def recentre(groups, card_w=3.99, gap=0.18, span_left=0.50, span_w=12.33):
    """Spread the surviving cards evenly instead of leaving a hole."""
    n = len(groups)
    if not n:
        return
    total = n * card_w + (n - 1) * gap
    x = span_left + (span_w - total) / 2.0
    for g in groups:
        cur = min(IN(sh.left) for sh in g)
        shift(g, dx=x - cur)
        x += card_w + gap


def keep_cards(slide, keep_titles, ytop=1.9, ybot=6.6):
    """Keep only the section cards whose text matches, then recentre."""
    kept, dropped = [], []
    for col in columns(slide, ytop, ybot):
        text = ' '.join(sh.text_frame.text for sh in col if sh.has_text_frame)
        (kept if any(t in text for t in keep_titles) else dropped).append(col)
    for g in dropped:
        drop_group(g)
    recentre(kept)
    return len(kept)


def keep_blocks(slide, keep_idx, ytop=1.5, ybot=5.2, pitch=0.75):
    """Keep objective rows by index, sliding survivors up into the gap."""
    allrows = rows(slide, ytop, ybot)
    for i, g in enumerate(allrows):
        if i not in keep_idx:
            drop_group(g)
    for new_pos, i in enumerate(sorted(keep_idx)):
        shift(allrows[i], dy=-(i - new_pos) * pitch)
    return len(keep_idx)


# --------------------------------------------------------------------------
# Content
# --------------------------------------------------------------------------
EK = {'TEACHER': True, 'STUDENT': False}

C = {
 '3.3': {
  'title': "Don't Let One Open Door Sink the Ship",
  'd1_sub': ('Protecting Networks: Segmentation - the three techniques that carve '
             'one flat network into separate compartments: the screened subnet, '
             'subnetting, and VLANs.'),
  'd2_sub': ('Protecting Networks: Segmentation - why isolating a network into '
             'segments reduces what an attacker can reach, and how each benefit '
             'traces back to the technique that produces it.'),
  'd1_keep_blocks': [0, 1],
  'd2_keep_blocks': [2],
  'd1_cards': ['Three Techniques'],
  'd2_cards': ['Why Segmentation Increases Security', 'Enrichment'],
  'd1_notes_intro': ("Your Day 1 Guided Notes have {n} section{s} that match "
                     "today's lesson. Green = CB required."),
  'd2_notes_intro': ("Your Day 2 Guided Notes have {n} section{s} that match "
                     "today's lesson. Green = CB required. Blue = enrichment."),
  'd2_bell': ("The same company has rebuilt its network. The public web server now "
              "sits in a screened subnet, the finance laptops are on their own "
              "VLAN, and the rest of the office is on a separate subnet. An "
              "employee opens the same malicious attachment on an office laptop."),
  'd2_task': ("Write 3-5 sentences: what can the malware reach now, what is out of "
              "its reach, and which single change did the most to contain it? Name "
              "the technique."),
  'learned': {
   'TEACHER': [
    "•  Named the three techniques that segment a network: the screened subnet or DMZ (EK 3.3.A.1), subnetting the address space (EK 3.3.A.2), and VLANs on a shared switch (EK 3.3.A.3)",
    "•  Told a physical division of the address space apart from a logical division of a single switch",
    "•  Settled the misconception that VLANs need their own separate hardware"],
   'STUDENT': [
    "•  Named the three techniques that segment a network: the screened subnet or DMZ, subnetting the address space, and VLANs on a shared switch",
    "•  Told a physical division of the address space apart from a logical division of a single switch",
    "•  Settled the misconception that VLANs need their own separate hardware"]},
  'next': {
   'TEACHER': [
    "•  Day 2: why segmentation actually makes a network safer (LO 3.3.B)",
    "•  Containment, different policies for higher- and lower-security zones, and port security against MAC flooding",
    "•  Matching each security benefit back to the technique and the EK that produce it"],
   'STUDENT': [
    "•  Day 2: why segmentation actually makes a network safer",
    "•  Containment, different policies for higher- and lower-security zones, and port security against MAC flooding",
    "•  Matching each security benefit back to the technique that produces it"]},
  'teaser': ("TEASER:  Today you carved the network up. Tomorrow the same infected "
             "laptop from the bell ringer comes back, so be ready to say exactly "
             "what it can still reach and which of the three techniques did the "
             "most to stop it."),
  'd1_notes': ("Day 1 of 2. Today is the three techniques only: screened subnet, "
               "subnetting, VLANs. Do not go into why segmentation works, that is "
               "tomorrow. Watch for students who think VLANs need separate "
               "switches; slide 9 exists for exactly that."),
  'd2_notes': ("Day 2 of 2. Yesterday named the three techniques; today is why they "
               "make a network safer. Open by calling back to the bell ringer, then "
               "run section 02. Section numbering continues from Day 1 on purpose."),
 },
 '3.4': {
  'title': 'The Bouncer at Every Door',
  'd1_sub': ('Protecting Networks: Firewalls - what a firewall is, the three CED '
             'types, and how an ordered access control list decides what gets '
             'through.'),
  'd2_sub': ('Protecting Networks: Firewalls - where firewalls belong on a '
             'segmented network, and how to turn written requirements into '
             'correctly ordered rules.'),
  'd1_keep_blocks': [0, 1],
  'd2_keep_blocks': [2, 3],
  'd1_cards': ['Firewall Types', 'The Access Control List'],
  'd2_cards': ['Firewall Placement', 'Configuring'],
  'd1_notes_intro': ("Your Day 1 Guided Notes have {n} section{s} that match "
                     "today's lesson. Green = CB required."),
  'd2_notes_intro': ("Your Day 2 Guided Notes have {n} section{s} that match "
                     "today's lesson. Green = CB required."),
  'd2_bell': ("A school network connects to the internet at two separate points. "
              "Inside, the student wifi, the staff subnet, and a server segment "
              "holding the grade database are all separated from one another. The "
              "district has bought one firewall and put it at one of the two "
              "internet connections."),
  'd2_task': ("Write 3-5 sentences: name every place this network still needs a "
              "firewall and say why. Should the grade database segment run the same "
              "rules as the student wifi? Explain."),
  'learned': {
   'TEACHER': [
    "•  Defined a firewall and named the three CED types, stateless, stateful and next-generation, and what each one inspects (EK 3.4.A.1-A.4)",
    "•  Read an access control list rule by rule and applied first match wins (EK 3.4.B.1-B.3)",
    "•  Settled the misconception that a stateful firewall and a next-generation firewall are the same thing"],
   'STUDENT': [
    "•  Defined a firewall and named the three types, stateless, stateful and next-generation, and what each one inspects",
    "•  Read an access control list rule by rule and applied first match wins",
    "•  Settled the misconception that a stateful firewall and a next-generation firewall are the same thing"]},
  'next': {
   'TEACHER': [
    "•  Day 2: where firewalls belong, at every segment and every internet ingress and egress point (LO 3.4.C)",
    "•  Writing rules from requirements and ordering them correctly (LO 3.4.D)",
    "•  Why two segments on one network can run different security levels"],
   'STUDENT': [
    "•  Day 2: where firewalls belong, at every segment and every internet ingress and egress point",
    "•  Writing rules from requirements and ordering them correctly",
    "•  Why two segments on one network can run different security levels"]},
  'teaser': ("TEASER:  You proved that the order of two rules changes what gets "
             "through. Tomorrow you write the rules yourself, from requirements, "
             "and find out how easily an ACL can permit exactly what you meant to "
             "block."),
  'd1_notes': ("Day 1 of 2. Firewall types and the ACL only. Placement and rule "
               "authoring are Day 2. The bell ringer already sets up first match "
               "wins, so refer back to it on slide 12 rather than re-explaining."),
  'd2_notes': ("Day 2 of 2. Yesterday was what a firewall is and how an ACL is "
               "read; today is where firewalls go and how to write the rules. "
               "Sections continue at 03, they are not renumbered."),
 },
}


def author_day1(lesson, variant, work):
    c = C[lesson]
    prs = Presentation(f'{work}/cut/{lesson}_D1_{variant}.pptx')
    src = Presentation(f'{work}/{lesson}_src_{variant}.pptx')
    tpl = Presentation(f'{work}/tpl_1.1_D1_{variant}.pptx')
    s1 = prs.slides[0]

    de.set_text(find(s1, 'DAY 1 OF'), 'DAY 1 OF 2')
    de.set_text(find(s1, f'Unit 3'), f'Unit 3 · Lesson {lesson} · Day 1 of 2')
    de.set_text(find(s1, 'Protecting Networks'), c['d1_sub'])
    s1.notes_slide.notes_text_frame.text = c['d1_notes']

    obj = prs.slides[1]
    keep_blocks(obj, c['d1_keep_blocks'])

    gn = find_slide(prs, 'GUIDED NOTES')
    n = keep_cards(gn, c['d1_cards'])
    de.set_text(find(gn, 'Guided Notes have'),
                c['d1_notes_intro'].format(n=n, s='' if n == 1 else 's'))

    # 3.4 Day 1 has no practice of its own: its only Stop and Think is slide 20,
    # which lands in Day 2. Copy the two questions that are scoped to A and B.
    if lesson == '3.4':
        # Located by content, not index: the STUDENT deck is missing the
        # teacher-only enrichment slide, so its Stop and Think sits one earlier.
        st = de.clone_slide(prs, find_slide(src, 'Stop and Think'))
        de.set_text(find(st, 'Complete on your'),
                    'Both questions map to LO 3.4.A and LO 3.4.B. Complete on your '
                    'Guided Notes packet.')
        qrows = rows(st, 1.6, 6.6)
        drop_group(qrows[2])
        st.notes_slide.notes_text_frame.text = (
            'Five minutes, independent. Question 1 is the type table, question 2 is '
            'first match wins. Placement is tomorrow, so do not let the discussion '
            'run to where firewalls belong.')

    close = de.clone_slide(prs, tpl.slides[len(tpl.slides) - 1])
    de.set_text(find(close, 'END OF DAY'), 'END OF DAY 1 · DAY 1 OF 2')
    de.set_text(find(close, 'End of Day'), 'End of Day 1')
    # The two bullet columns are the wide boxes on the same row, left and
    # right. Matching on their text would tie this to Lesson 1.1's wording.
    cols = sorted((sh for sh in close.shapes
                   if sh.has_text_frame and 2.3 <= IN(sh.top) <= 3.0
                   and IN(sh.width) > 4.0),
                  key=lambda sh: IN(sh.left))
    if len(cols) != 2:
        raise SystemExit(f'day-close: expected 2 bullet columns, found {len(cols)}')
    de.set_bullets(cols[0], c['learned'][variant])
    de.set_bullets(cols[1], c['next'][variant])
    de.set_text(find(close, 'TEASER'), c['teaser'])
    close.notes_slide.notes_text_frame.text = (
        f"Close Day 1 here. Recap the {'techniques' if lesson=='3.3' else 'types and the ACL'}, "
        f"then read the teaser aloud and stop. Day 2 opens by answering it.")

    de.renumber_footers(prs, lesson)
    out = f'{work}/final/{lesson}_Day1_Deck_{variant}.pptx'
    prs.save(out)
    return out, len(prs.slides)


def find_slide(prs, needle):
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and needle in sh.text_frame.text:
                return s
    raise KeyError(needle)


def author_day2(lesson, variant, work):
    c = C[lesson]
    prs = Presentation(f'{work}/cut/{lesson}_D2_{variant}.pptx')
    src = Presentation(f'{work}/{lesson}_src_{variant}.pptx')

    # Front matter, cloned from the lesson's own Day 1 slides so the EDITION
    # label, theme and footer shape all come along correct by construction.
    title = de.clone_slide(prs, src.slides[0])
    de.set_text(find(title, 'DAY 1 OF'), 'DAY 2 OF 2')
    de.set_text(find(title, 'Unit 3'), f'Unit 3 · Lesson {lesson} · Day 2 of 2')
    de.set_text(find(title, 'Protecting Networks'), c['d2_sub'])
    title.notes_slide.notes_text_frame.text = c['d2_notes']

    obj = de.clone_slide(prs, src.slides[1])
    keep_blocks(obj, c['d2_keep_blocks'])
    obj.notes_slide.notes_text_frame.text = (
        "Today's objectives only. Yesterday's are done and are not repeated here.")

    bell = de.clone_slide(prs, src.slides[2])
    de.set_text(find(bell, 'A small company') if lesson == '3.3'
                else find(bell, 'A network administrator'), c['d2_bell'])
    de.set_text(find(bell, 'Write 3'), c['d2_task'])
    bell.notes_slide.notes_text_frame.text = (
        'Five minutes. This deliberately calls back to Day 1 so the two days read '
        'as one lesson rather than two unrelated decks.')

    gn = de.clone_slide(prs, src.slides[3])
    n = keep_cards(gn, c['d2_cards'])
    de.set_text(find(gn, 'Guided Notes have'),
                c['d2_notes_intro'].format(n=n, s='' if n == 1 else 's'))
    gn.notes_slide.notes_text_frame.text = (
        'Section numbers continue from Day 1. They are not renumbered, which is '
        'the same convention Unit 1 uses.')

    # The four new slides were appended; move them to the front, in order.
    last = len(prs.slides) - 1
    for offset, _ in enumerate([title, obj, bell, gn]):
        de.move_slide(prs, last - 3 + offset, offset)

    de.renumber_footers(prs, lesson)
    out = f'{work}/final/{lesson}_Day2_Deck_{variant}.pptx'
    prs.save(out)
    return out, len(prs.slides)


if __name__ == '__main__':
    work = sys.argv[1]
    os.makedirs(f'{work}/final', exist_ok=True)
    for lesson in ('3.3', '3.4'):
        for variant in ('TEACHER', 'STUDENT'):
            p1, n1 = author_day1(lesson, variant, work)
            p2, n2 = author_day2(lesson, variant, work)
            print(f'{lesson} {variant:8s} Day1 {n1:2d} slides   Day2 {n2:2d} slides')
