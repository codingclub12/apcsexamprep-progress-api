#!/usr/bin/env python3
"""Range-split one AP Cybersecurity lesson deck into per-day decks.

    python3 split_deck.py <in.pptx> <out.pptx> <first> <last>

Keeps slides [first..last] (1-based, inclusive) and drops the rest, leaving
masters, layouts, theme, media and speaker notes untouched.

WHY A SCRIPT, when docs/cyber-unit3-tier1-split-spec.md says not to bother.
That line was written for a human with PowerPoint open, and for four decks it
was right. Units 3-5 are 15 lessons and 30 source decks going to about 130, and
no session of Claude Code has PowerPoint. The mechanical cut is also the part
where a human miscounts; the authoring the spec cares about is still authoring.

WHAT THIS DELIBERATELY DOES NOT DO. It does not renumber footers, rewrite the
DAY n OF m badge, rescope the objectives or guided-notes slides, or touch the
title slide's speaker notes. Every one of those is per-day authoring and every
one is visible to a teacher; see the spec's "Edits that apply to every day
deck". A deck that has only been through this script is not shippable, and
`audit_deck.py` in this directory is what says so out loud rather than letting
a half-done deck look finished.

python-pptx has no public delete-slide API. Dropping the sldId entry and its
relationship is the documented workaround: the slide part stops being
referenced and is not written out on save.
"""
import sys
from pptx import Presentation


def slide_range(path_in, path_out, first, last):
    prs = Presentation(path_in)
    total = len(prs.slides)
    if not (1 <= first <= last <= total):
        raise SystemExit(
            f'range {first}-{last} is not inside 1-{total} for {path_in}')

    id_list = prs.slides._sldIdLst
    entries = list(id_list)
    # Drop from the end so earlier indices stay valid as we mutate.
    for idx in range(total - 1, -1, -1):
        if first <= idx + 1 <= last:
            continue
        entry = entries[idx]
        prs.part.drop_rel(entry.rId)
        id_list.remove(entry)

    prs.save(path_out)
    kept = len(Presentation(path_out).slides)
    expected = last - first + 1
    if kept != expected:
        raise SystemExit(f'expected {expected} slides, wrote {kept}')
    return total, kept


if __name__ == '__main__':
    if len(sys.argv) != 5:
        raise SystemExit(__doc__)
    src, dst, a, b = sys.argv[1], sys.argv[2], int(sys.argv[3]), int(sys.argv[4])
    was, now = slide_range(src, dst, a, b)
    print(f'{src} [{was} slides] -> {dst} [{now} slides, kept {a}-{b}]')
